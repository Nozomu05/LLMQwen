"""Search Service — entry point for the RAG pipeline.

Receives: POST /query  {"query", "documents", "temperature", "agent_id", "history"}
Flow:     language detection → translation → embedding → Chroma vector search
          (filtered to the agent's docs/ folder) → optional SearXNG web
          search (if the agent has web_search enabled) → POST /rerank
          (Reranker Service) → POST /generate (LLM Service)
Returns:  {"answer", "model", "sources"}
"""
import os
import re
import time
import warnings
from io import BytesIO
from pathlib import Path
from typing import Optional

import requests
from dotenv import load_dotenv
from fastapi import FastAPI, File, HTTPException, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse
from pydantic import BaseModel

from services.agents_config import DEFAULT_AGENT_ID, get_agent, list_agents, reload_agent, reload_all_agents

warnings.filterwarnings(
    "ignore",
    message=".*multilingual-e5-large now uses mean pooling.*",
    category=UserWarning,
)

load_dotenv(override=True)

app = FastAPI(title="Search Service")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

_EMBEDDINGS_CACHE: dict = {}
_VECTORSTORE_CACHE: dict = {}

MAX_EXTRACT_CHARS = 20_000

_LANG_INSTRUCTIONS = {
    "fr": "Réponds en français.",
    "en": "Respond in English.",
}

# langdetect is a statistical, multi-language classifier and is unreliable on
# short text — e.g. it misreads "courses"/"requirements" as French, and
# "Qui est Marius Preda ?" as Romanian. Since we only ever need a binary
# French/English call for the response language, a small dedicated heuristic
# (accented characters + common French words) is more robust than trusting
# langdetect's full N-language guess for short queries.
_FRENCH_CHARS = set("àâçéèêëîïôùûüœ")
_FRENCH_WORDS = {
    "le", "la", "les", "un", "une", "des", "du", "de", "et", "est", "es", "suis",
    "qui", "que", "quoi", "quel", "quels", "quelle", "quelles", "sont", "avez",
    "peux", "peut", "pouvez", "voudrais", "aimerais", "cherche",
    "comment", "pourquoi", "combien", "quand", "bonjour", "merci",
    "vous", "nous", "je", "tu", "il", "elle", "avec", "pour", "dans", "sur",
    "sans", "mais", "donc", "alors", "aussi",
    "frais", "cours", "inscription", "bourse", "bourses", "logement",
    "logements", "stage", "stages", "recherche", "semestre", "contact",
    "contacter", "campus", "scolarite",
}


def _looks_french(text: str) -> bool:
    lowered = text.lower()
    if any(ch in _FRENCH_CHARS for ch in lowered):
        return True
    words = re.findall(r"[a-zà-ÿ']+", lowered)
    return any(w in _FRENCH_WORDS for w in words)


# ---------------------------------------------------------------------------
# Pydantic models
# ---------------------------------------------------------------------------

class Document(BaseModel):
    filename: str
    text: str


class QueryRequest(BaseModel):
    query: str
    documents: Optional[list[Document]] = None
    temperature: Optional[float] = None
    agent_id: str = DEFAULT_AGENT_ID
    # Prior conversation turns, [{"role": "user"|"assistant", "content": str}, ...].
    # Passed straight through to the LLM service for context — deliberately
    # NOT used for vector/web search, which stay scoped to the current query
    # alone (concatenating history into the search text is what broke
    # retrieval before: it searched for the whole conversation, not the question).
    history: Optional[list[dict]] = None


class ReloadAgentsRequest(BaseModel):
    agent_id: Optional[str] = None


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _get_embeddings():
    from langchain_community.embeddings import FastEmbedEmbeddings
    from langchain_huggingface import HuggingFaceEmbeddings

    provider = os.environ["EMBEDDING_PROVIDER_ML"].lower()
    model = os.environ["EMBEDDING_MODEL_ML"]
    key = (provider, model)
    if key in _EMBEDDINGS_CACHE:
        return _EMBEDDINGS_CACHE[key]

    normalize = os.environ["EMBEDDING_NORMALIZE"].lower() == "true"
    if provider == "huggingface":
        emb = HuggingFaceEmbeddings(
            model_name=model,
            model_kwargs={"device": os.environ["EMBEDDING_DEVICE"]},
            encode_kwargs={"normalize_embeddings": normalize},
        )
    else:
        emb = FastEmbedEmbeddings(
            model_name=model,
            max_length=512,
            additional_kwargs={"providers": ["CPUExecutionProvider"]},
        )
    _EMBEDDINGS_CACHE[key] = emb
    return emb


def _detect_language(text: str) -> tuple[str, str]:
    """Return (lang_code, response_hint).

    The chatbot only ever answers in French or English: French if the
    question is detected as French, English for every other language. The
    response-language call uses _looks_french() rather than langdetect,
    which is unreliable on short queries (see _looks_french's docstring).

    `lang_code` still comes from langdetect and reflects the actual detected
    language — used elsewhere to decide whether to translate the retrieval
    query to English. Being wrong there is low-stakes (mistranslating
    English to English is a no-op), unlike being wrong about which language
    to respond in, so it doesn't need the same override.
    """
    if _looks_french(text):
        return "fr", f" ({_LANG_INSTRUCTIONS['fr']})"
    try:
        from langdetect import detect
        lang_code = detect(text)
    except Exception:
        return "en", ""
    if lang_code == "en":
        return "en", ""
    return lang_code, f" ({_LANG_INSTRUCTIONS['en']})"


def _translate_to_english(text: str) -> str:
    try:
        from deep_translator import GoogleTranslator
        return GoogleTranslator(source="auto", target="en").translate(text)
    except Exception:
        return text


def _get_vectorstore():
    from langchain_chroma import Chroma

    chroma_dir = Path(os.environ["CHROMA_DIR"]).resolve()
    if not chroma_dir.exists():
        raise FileNotFoundError(
            f"No vector store at {chroma_dir}. Run ingestion first: python rag/ingest.py"
        )
    vs_key = str(chroma_dir)
    if vs_key not in _VECTORSTORE_CACHE:
        _VECTORSTORE_CACHE[vs_key] = Chroma(
            persist_directory=str(chroma_dir),
            embedding_function=_get_embeddings(),
        )
    return _VECTORSTORE_CACHE[vs_key]


def _extract_text_from_bytes(filename: str, file_bytes: bytes) -> str:
    """Convert an uploaded office/PDF file to plain text for use as a query document."""
    ext = Path(filename).suffix.lower()
    try:
        if ext == ".pdf":
            import fitz
            doc = fitz.open(stream=file_bytes, filetype="pdf")
            text = "\n".join(page.get_text() for page in doc)
        elif ext == ".docx":
            from docx import Document as DocxDoc
            doc = DocxDoc(BytesIO(file_bytes))
            text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        elif ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(BytesIO(file_bytes))
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        parts.append(shape.text)
            text = "\n".join(parts)
        elif ext in (".xlsx", ".xls", ".ods"):
            import openpyxl
            wb = openpyxl.load_workbook(BytesIO(file_bytes), read_only=True, data_only=True)
            rows = []
            for ws in wb.worksheets:
                rows.append(f"[Sheet: {ws.title}]")
                for row in ws.iter_rows(values_only=True):
                    row_str = "\t".join(str(c) if c is not None else "" for c in row)
                    if row_str.strip():
                        rows.append(row_str)
            text = "\n".join(rows)
        else:
            text = file_bytes.decode("utf-8", errors="replace")
    except Exception as e:
        text = f"[Error extracting text from {filename}: {e}]"

    if len(text) > MAX_EXTRACT_CHARS:
        text = text[:MAX_EXTRACT_CHARS] + f"\n\n[... truncated at {MAX_EXTRACT_CHARS:,} characters]"
    return text


def _vector_search(retrieval_query: str, agent_folder: str) -> list[dict]:
    """Embed the retrieval query and search the Chroma index, restricted to
    chunks tagged with `agent_folder` (the docs/ subfolder the agent owns)."""
    k_total = int(os.environ["RETRIEVAL_CHUNKS"])
    vs = _get_vectorstore()

    # Instruct-style models need a prefix on the query side
    embedding_model = os.environ["EMBEDDING_MODEL_ML"]
    if "instruct" in embedding_model.lower():
        retrieval_query = (
            f"Instruct: Retrieve relevant passages that answer the question\n"
            f"Query: {retrieval_query}"
        )

    docs = vs.as_retriever(
        search_kwargs={"k": k_total, "filter": {"agent": agent_folder}}
    ).invoke(retrieval_query)
    return [
        {"content": d.page_content, "source": d.metadata.get("source", "unknown"), "origin": "document"}
        for d in docs
    ]


_WEB_SEARCH_CACHE: dict = {}
_WEB_SEARCH_CACHE_TTL = 300  # seconds


def _web_search(query: str, context: str = "") -> list[dict]:
    """Search the web via a local SearXNG instance and return results shaped
    like vector-search chunks (content + source URL), so they flow through
    the same reranking and [Source: ...] citation path as document chunks.

    `context` is an agent-specific scoping phrase (e.g. "Télécom SudParis",
    from the agent's web_search_context in the database) appended to the
    query so results stay on-topic instead of drifting to whatever else the
    bare user query happens to match on the open web.

    SearXNG aggregates multiple engines (Google, Startpage, DuckDuckGo, ...)
    per query, so one engine having a bad moment doesn't sink the whole
    result set the way a single-source scraper can. Results are still
    cached for a few minutes so repeated identical questions don't refetch.

    This always fails soft (empty list, logged) instead of raising, so an
    unreachable SearXNG instance never breaks the whole /query request for
    agents that don't strictly need it.
    """
    scoped_query = f"{query} {context}".strip() if context else query

    cached = _WEB_SEARCH_CACHE.get(scoped_query)
    if cached and time.time() - cached[0] < _WEB_SEARCH_CACHE_TTL:
        return cached[1]

    try:
        searxng_url = os.environ["SEARXNG_URL"]
        n_results = int(os.environ.get("WEB_SEARCH_RESULTS", "5"))
        resp = requests.get(
            f"{searxng_url}/search",
            params={"q": scoped_query, "format": "json"},
            timeout=15,
        )
        resp.raise_for_status()
        results = resp.json().get("results", [])[:n_results]
        chunks = [
            {
                "content": f"{r.get('title', '')}\n{r.get('content', '')}".strip(),
                "source": r["url"],
                "origin": "web",
            }
            for r in results
            if r.get("url")
        ]
        _WEB_SEARCH_CACHE[scoped_query] = (time.time(), chunks)
        return chunks
    except Exception as exc:
        print(f"[search] web search failed: {exc}")
        return []


# ---------------------------------------------------------------------------
# Startup
# ---------------------------------------------------------------------------

@app.on_event("startup")
def warmup():
    _get_embeddings()
    try:
        _get_vectorstore()
    except FileNotFoundError:
        pass  # chroma index not ingested yet


# ---------------------------------------------------------------------------
# Endpoints
# ---------------------------------------------------------------------------

@app.get("/health")
def health():
    return {"status": "ok"}


@app.get("/agents")
def agents():
    """List every agent (id, name, profession, department, web_search) for
    UIs to populate an agent selector dynamically."""
    return {"agents": list_agents()}


@app.post("/agents/reload")
def reload_agents(req: ReloadAgentsRequest):
    """Drop cached agent folder mapping(s) so the next request re-reads them
    from the database, and cascade the same reload to the LLM service so a
    single call refreshes both cache layers (folder/web_search here, persona
    system_prompt there)."""
    if req.agent_id:
        reload_agent(req.agent_id)
    else:
        reload_all_agents()
    try:
        requests.post(
            f"{os.environ['LLM_SERVICE_URL']}/agents/reload",
            json={"agent_id": req.agent_id},
            timeout=10,
        )
    except Exception as exc:
        print(f"[search] failed to cascade agent reload to LLM service: {exc}")
    return {"status": "ok"}


@app.post("/extract")
async def extract(file: UploadFile = File(...)):
    """Extract plain text from an uploaded PDF/DOCX/PPTX/XLSX file for use as a query document."""
    file_bytes = await file.read()
    text = _extract_text_from_bytes(file.filename or "document", file_bytes)
    return {"filename": file.filename, "text": text}


@app.post("/query")
def query(req: QueryRequest):
    """Full RAG entry point: search → rerank → generate."""
    t_start = time.time()
    print(f"[search] /query: temperature={req.temperature!r} agent_id={req.agent_id!r} query={req.query!r}")

    reranker_url = os.environ["RERANKER_SERVICE_URL"]
    llm_url = os.environ["LLM_SERVICE_URL"]

    agent = get_agent(req.agent_id)
    agent_folder = agent["folder"]
    lang_code, lang_hint = _detect_language(req.query)

    # Determine the retrieval query and the reranking query
    if os.environ["USE_HYDE"].lower() == "true":
        # HyDE: generate a hypothetical answer to use as the retrieval query.
        # The hypothesis is also used as the rerank query (matches original behaviour).
        try:
            resp = requests.post(
                f"{llm_url}/hyde",
                json={"query": req.query},
                timeout=60,
            )
            resp.raise_for_status()
            hypothesis = resp.json().get("hypothesis", req.query)
        except Exception:
            hypothesis = req.query
        retrieval_query = hypothesis
        rerank_query = hypothesis
    else:
        # Translate to English for better embedding similarity (all docs are in English).
        # Use the original (possibly non-English) query for reranking — the reranker is multilingual.
        if lang_code != "en":
            retrieval_query = _translate_to_english(req.query)
        else:
            retrieval_query = req.query
        rerank_query = req.query

    try:
        t_search = time.time()
        chunks = _vector_search(retrieval_query, agent_folder)
        print(f"[search] vector search: {time.time() - t_search:.2f}s — {len(chunks)} chunks retrieved")
    except FileNotFoundError as exc:
        raise HTTPException(status_code=503, detail=str(exc))

    if agent["web_search"]:
        t_web = time.time()
        web_chunks = _web_search(req.query, agent["web_search_context"])
        print(f"[search] web search: {time.time() - t_web:.2f}s — {len(web_chunks)} results")
        chunks = chunks + web_chunks

    extra_docs = None
    if req.documents:
        extra_docs = [{"filename": d.filename, "text": d.text} for d in req.documents]

    rerank_payload = {
        "query": req.query,
        "chunks": chunks,
        "rerank_query": rerank_query,
        "lang_hint": lang_hint,
        "extra_docs": extra_docs,
        "temperature": req.temperature,
        "agent_id": req.agent_id,
        "history": req.history,
    }

    def proxy_stream():
        t_reranker = time.time()
        try:
            with requests.post(
                f"{reranker_url}/rerank",
                json=rerank_payload,
                stream=True,
                timeout=180,
            ) as resp:
                if resp.status_code != 200:
                    import json as _json
                    yield f"data: {_json.dumps({'type': 'error', 'error': f'Reranker service returned HTTP {resp.status_code}'})}\n\n".encode()
                else:
                    for raw in resp.iter_content(chunk_size=None):
                        yield raw
        except Exception as exc:
            import json as _json
            yield f"data: {_json.dumps({'type': 'error', 'error': f'Reranker service error: {exc}'})}\n\n".encode()
        print(f"[search] rerank+generate: {time.time() - t_reranker:.2f}s")
        print(f"[search] total: {time.time() - t_start:.2f}s")

    return StreamingResponse(proxy_stream(), media_type="text/event-stream")
